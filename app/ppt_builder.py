from __future__ import annotations

from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.models import GeneratedSlideContent, PPTGenerationInput, ReferenceContext


# 主题色
_TITLE_COLOR = RGBColor(0x1A, 0x47, 0x8A)      # 深蓝
_BODY_COLOR = RGBColor(0x33, 0x33, 0x33)        # 深灰
_ACCENT_COLOR = RGBColor(0x2E, 0x86, 0xAB)      # 亮蓝

# 页面尺寸（16:9）
_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)


def _set_font(run, size_pt: int, color: RGBColor, bold: bool = False):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "微软雅黑"


def _add_picture_if_exists(slide, image_path: str, left, top, width=None, height=None) -> bool:
    if not image_path:
        return False
    path = Path(image_path)
    if not path.exists():
        return False
    if width is not None and height is not None:
        slide.shapes.add_picture(str(path), left, top, width, height)
    elif width is not None:
        slide.shapes.add_picture(str(path), left, top, width=width)
    elif height is not None:
        slide.shapes.add_picture(str(path), left, top, height=height)
    else:
        slide.shapes.add_picture(str(path), left, top)
    return True


def _add_reference_logo(slide, ref_context: ReferenceContext | None, is_title: bool = False) -> None:
    if not ref_context or not ref_context.logo_path:
        return
    if is_title:
        _add_picture_if_exists(
            slide,
            ref_context.logo_path,
            Inches(11.2),
            Inches(0.35),
            width=Inches(1.1),
        )
    else:
        _add_picture_if_exists(
            slide,
            ref_context.logo_path,
            Inches(10.55),
            Inches(0.18),
            width=Inches(0.75),
        )


def _add_title_slide(prs: Presentation, input_data: PPTGenerationInput, ref_context: ReferenceContext | None = None):
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    w = _SLIDE_WIDTH
    h = _SLIDE_HEIGHT

    if ref_context and ref_context.background_image_path:
        _add_picture_if_exists(
            slide,
            ref_context.background_image_path,
            0,
            0,
            width=w,
            height=h,
        )

    # 顶部色块
    top_bar = slide.shapes.add_shape(
        1, 0, 0, w, Inches(0.15)  # MSO_SHAPE.RECTANGLE = 1
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = _TITLE_COLOR
    top_bar.line.fill.background()

    _add_reference_logo(slide, ref_context, is_title=True)

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = input_data.topic
    _set_font(run, 36, _TITLE_COLOR, bold=True)

    # 副标题 / 用途
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"—— {input_data.purpose} ——"
    company = (ref_context.company_name if ref_context and ref_context.company_name else "")
    if company:
        txBox_company = slide.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10), Inches(0.5))
        tf_company = txBox_company.text_frame
        tf_company.word_wrap = True
        p_company = tf_company.paragraphs[0]
        p_company.alignment = PP_ALIGN.CENTER
        run_company = p_company.add_run()
        run_company.text = company
        _set_font(run_company, 16, _BODY_COLOR)
    _set_font(run2, 18, _ACCENT_COLOR)

    # 底部信息
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = f"目标受众：{input_data.audience}   风格：{input_data.style}"
    _set_font(run3, 14, _BODY_COLOR)


def _add_content_slide(
    prs: Presentation,
    content: GeneratedSlideContent,
    slide_number: int,
    total: int,
    ref_context: ReferenceContext | None = None,
):
    """创建正文页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    w = _SLIDE_WIDTH

    # 顶部色条
    top_bar = slide.shapes.add_shape(1, 0, 0, w, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = _ACCENT_COLOR
    top_bar.line.fill.background()

    _add_reference_logo(slide, ref_context)

    # 页码
    txPage = slide.shapes.add_textbox(Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.4))
    tfp = txPage.text_frame
    runp = tfp.paragraphs[0].add_run()
    runp.text = f"{slide_number} / {total}"
    _set_font(runp, 12, RGBColor(0x99, 0x99, 0x99))

    # 标题
    txTitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8))
    tft = txTitle.text_frame
    tft.word_wrap = True
    runt = tft.paragraphs[0].add_run()
    runt.text = content.title
    _set_font(runt, 28, _TITLE_COLOR, bold=True)

    # 分隔线
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = _ACCENT_COLOR
    line.line.fill.background()

    # 正文要点
    txBody = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5))
    tfb = txBody.text_frame
    tfb.word_wrap = True

    for i, bullet in enumerate(content.bullets):
        if i == 0:
            p = tfb.paragraphs[0]
        else:
            p = tfb.add_paragraph()

        p.space_after = Pt(8)
        p.level = 0

        run = p.add_run()
        # 要点标记
        run.text = f"  ▌  {bullet}"
        body_color = _BODY_COLOR
        if ref_context and ref_context.theme_colors and len(ref_context.theme_colors) >= 2:
            try:
                body_color = RGBColor(int(ref_context.theme_colors[1], 16))
            except Exception:
                pass
        _set_font(run, 18, body_color)

    # speaker_hint 在最底部
    if content.speaker_hint:
        txHint = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6))
        tfh = txHint.text_frame
        tfh.word_wrap = True
        runh = tfh.paragraphs[0].add_run()
        runh.text = f"💡 {content.speaker_hint}"
        _set_font(runh, 12, RGBColor(0x88, 0x88, 0x88))
        runh.font.italic = True


def build_pptx(
    input_data: PPTGenerationInput,
    slides_content: Sequence[GeneratedSlideContent],
    output_path: str | Path,
    ref_context: ReferenceContext | None = None,
) -> Path:
    """从结构化内容生成 .pptx 文件"""
    output_path = Path(output_path)
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT

    # 去除所有默认占位符（从 blank layout 来的空框）
    # 我们用 blank layout (index 6)，没有占位符

    total = len(slides_content)
    if total == 0:
        raise ValueError("没有可生成的页面内容。")

    # 封面页
    _add_title_slide(prs, input_data, ref_context)

    # 正文页
    for i, content in enumerate(slides_content):
        _add_content_slide(prs, content, i + 1, total, ref_context)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path.resolve()

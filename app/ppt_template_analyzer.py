
from __future__ import annotations

from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.models import ReferenceContext


def _extract_doc_properties(prs: Presentation) -> dict:
    """提取文档属性的公司/作者信息"""
    info = {"company": "", "author": "", "title": ""}
    cp = prs.core_properties
    if hasattr(cp, 'company') and cp.company:
        info["company"] = cp.company.strip()
    if hasattr(cp, 'author') and cp.author:
        info["author"] = cp.author.strip()
    if hasattr(cp, 'title') and cp.title:
        info["title"] = cp.title.strip()
    return info


def _extract_theme_colors(prs: Presentation) -> list[str]:
    """从主题中提取主要配色"""
    colors = []
    try:
        theme_elem = prs.slide_masters[0].element
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_scheme = theme_elem.find(".//a:clrScheme", ns)
        if clr_scheme is not None:
            for child in clr_scheme:
                for sub in child:
                    val = sub.get("val") or sub.get("lastClr")
                    if val:
                        colors.append(val)
    except Exception:
        pass
    # 去重保留顺序
    seen = set()
    unique = []
    for c in colors:
        if c not in seen:
            seen.add(c)
            unique.append(c)
    return unique[:6]


def _extract_fonts_and_colors(prs: Presentation) -> tuple[list[str], list[str]]:
    """从幻灯片形状中提取实际使用的字体和颜色"""
    font_counter: Counter = Counter()
    color_counter: Counter = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        font_counter[run.font.name] += 1
                    try:
                        if run.font.color and run.font.color.rgb:
                            c = str(run.font.color.rgb)
                            color_counter[c] += 1
                    except Exception:
                        pass
    fonts = [f for f, _ in font_counter.most_common(3)]
    colors = [c for c, _ in color_counter.most_common(6)]
    return fonts, colors


def _extract_image_assets(
    prs: Presentation,
    asset_dir: str | Path,
) -> tuple[str, str, list[str], str]:
    """提取参考 PPT 图片，并粗略识别 logo 和背景图。"""
    asset_root = Path(asset_dir)
    asset_root.mkdir(parents=True, exist_ok=True)

    image_paths: list[str] = []
    logo_candidates: list[tuple[float, str]] = []
    background_candidates: list[tuple[float, str]] = []
    slide_area = int(prs.slide_width) * int(prs.slide_height)

    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            image = shape.image
            ext = image.ext or "png"
            image_path = asset_root / f"slide_{slide_index:03d}_image_{shape_index:02d}.{ext}"
            image_path.write_bytes(image.blob)
            resolved = str(image_path.resolve())
            image_paths.append(resolved)

            shape_area = int(shape.width) * int(shape.height)
            area_ratio = shape_area / slide_area if slide_area else 0
            near_edge = (
                shape.left < prs.slide_width * 0.2
                or shape.left > prs.slide_width * 0.7
                or shape.top < prs.slide_height * 0.2
                or shape.top > prs.slide_height * 0.7
            )
            if area_ratio <= 0.08 and near_edge:
                logo_candidates.append((area_ratio, resolved))
            if area_ratio >= 0.5:
                background_candidates.append((area_ratio, resolved))

    logo_path = logo_candidates[0][1] if logo_candidates else (image_paths[0] if image_paths else "")
    background_path = (
        sorted(background_candidates, reverse=True)[0][1]
        if background_candidates
        else ""
    )
    return str(asset_root.resolve()), logo_path, background_path, image_paths


def analyze_reference_pptx(
    pptx_path: str | Path,
    asset_dir: str | Path | None = None,
) -> ReferenceContext:
    """解析参考 PPT，提取页数、标题、公司名、配色、字体等风格信息"""
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if asset_dir is None:
        asset_dir = Path("temp") / "reference_assets" / pptx_path.stem
    
    # 提取幻灯片文本结构
    slides_summary: list[dict] = []
    for zero_based, slide in enumerate(prs.slides):
        texts: list[str] = []
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
            texts.append(title)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                t = shape.text.strip()
                if t != title:
                    texts.append(t)
        raw = "\n".join(texts).strip()
        body_raw = "\n".join(t for t in texts if t != title).strip()
        body_preview = body_raw[:120]
        if len(body_raw) > 120:
            body_preview += "..."
        slides_summary.append({
            "index": zero_based + 1,
            "title": title,
            "body_preview": body_preview,
        })
    
    # 提取公司/作者信息
    doc_info = _extract_doc_properties(prs)
    
    # 提取主题色
    theme_colors = _extract_theme_colors(prs)
    
    # 提取实际使用的字体和颜色
    fonts, shape_colors = _extract_fonts_and_colors(prs)
    
    # 合并颜色源
    all_colors = theme_colors + shape_colors
    seen = set()
    merged_colors = []
    for c in all_colors:
        if c not in seen:
            seen.add(c)
            merged_colors.append(c)
    merged_colors = merged_colors[:6]
    
    # 构建风格描述文本（供 LLM + 页面展示用）
    style_parts = []
    if doc_info["company"]:
        style_parts.append(f"公司：{doc_info['company']}")
    if fonts:
        style_parts.append(f"字体：{', '.join(fonts)}")
    if merged_colors:
        style_parts.append(f"主色：{' '.join('#' + c for c in merged_colors)}")
    asset_root, logo_path, background_path, image_paths = _extract_image_assets(
        prs, asset_dir
    )
    if logo_path:
        style_parts.append("已提取参考 Logo/图片素材")
    style_notes = "；".join(style_parts)
    
    return ReferenceContext(
        slide_count=len(slides_summary),
        slides_summary=slides_summary,
        company_name=doc_info["company"],
        theme_colors=merged_colors,
        font_name=fonts[0] if fonts else "",
        style_notes=style_notes,
        asset_dir=asset_root,
        logo_path=logo_path,
        background_image_path=background_path,
        image_paths=image_paths,
    )

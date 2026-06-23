from __future__ import annotations

from pathlib import Path
from typing import BinaryIO

from pptx import Presentation

from app.models import SlideContent


def _shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return str(shape.text).strip()


def parse_pptx(file: BinaryIO | str | Path) -> list[SlideContent]:
    presentation = Presentation(file)
    slides: list[SlideContent] = []

    for zero_based_index, slide in enumerate(presentation.slides):
        text_parts: list[str] = []
        title = ""

        if slide.shapes.title is not None:
            title = _shape_text(slide.shapes.title)
            if title:
                text_parts.append(title)

        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text:
                continue
            if text == title:
                continue
            text_parts.append(text)

        raw_text = "\n".join(text_parts).strip()
        body = "\n".join(part for part in text_parts if part != title).strip()
        slides.append(
            SlideContent(
                index=zero_based_index + 1,
                title=title,
                body=body,
                raw_text=raw_text,
            )
        )

    readable_slides = [slide for slide in slides if slide.raw_text]
    if not readable_slides:
        raise ValueError("PPT 中没有可提取的文本。")
    return readable_slides

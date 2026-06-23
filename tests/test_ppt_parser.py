from pathlib import Path

import pytest
from pptx import Presentation

from app.ppt_parser import parse_pptx


def _create_sample_pptx(path: Path) -> None:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "第一页标题"
    slide.placeholders[1].text = "第一页正文"
    presentation.save(path)


def test_parse_pptx_extracts_title_body_and_raw_text(tmp_path: Path):
    pptx_path = tmp_path / "sample.pptx"
    _create_sample_pptx(pptx_path)

    slides = parse_pptx(pptx_path)

    assert len(slides) == 1
    assert slides[0].index == 1
    assert slides[0].title == "第一页标题"
    assert slides[0].body == "第一页正文"
    assert slides[0].raw_text == "第一页标题\n第一页正文"


def test_parse_pptx_raises_when_no_text_can_be_extracted(tmp_path: Path):
    pptx_path = tmp_path / "empty.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(pptx_path)

    with pytest.raises(ValueError, match="没有可提取的文本"):
        parse_pptx(pptx_path)

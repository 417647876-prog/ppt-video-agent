from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.models import PPTGenerationInput, PPTOutlineItem, ReferenceContext
from app.ppt_template_analyzer import analyze_reference_pptx


def _make_minimal_pptx(tmp_path: Path, slides_data: list[tuple[str, str]]) -> Path:
    """创建一个最小测试 PPTX 文件"""
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    for title, body in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        # 添加标题
        txBox = slide.shapes.add_textbox(0, 0, 100, 20)
        txBox.text_frame.paragraphs[0].add_run()
        txBox.text_frame.paragraphs[0].runs[0].text = title
        # 添加正文
        txBody = slide.shapes.add_textbox(0, 30, 100, 100)
        txBody.text_frame.paragraphs[0].add_run()
        txBody.text_frame.paragraphs[0].runs[0].text = body
    prs.save(str(pptx_path))
    return pptx_path


def test_analyze_reference_pptx_returns_correct_slide_count(tmp_path: Path):
    pptx_path = _make_minimal_pptx(
        tmp_path,
        [("封面", "欢迎"), ("正文", "核心内容"), ("结尾", "谢谢")],
    )
    context = analyze_reference_pptx(pptx_path)
    assert context.slide_count == 3
    assert len(context.slides_summary) == 3


def test_analyze_reference_pptx_extracts_titles(tmp_path: Path):
    pptx_path = _make_minimal_pptx(
        tmp_path,
        [("第一页标题", "内容一"), ("第二页标题", "内容二")],
    )
    context = analyze_reference_pptx(pptx_path)
    assert "第一页标题" in context.slides_summary[0]["body_preview"]
    assert "第二页标题" in context.slides_summary[1]["body_preview"]


def test_analyze_reference_pptx_body_preview_not_empty(tmp_path: Path):
    pptx_path = _make_minimal_pptx(tmp_path, [("标题", "这是一段正文内容")])
    context = analyze_reference_pptx(pptx_path)
    assert len(context.slides_summary[0]["body_preview"]) > 0
    assert "正文内容" in context.slides_summary[0]["body_preview"]


def test_analyze_reference_pptx_extracts_image_assets(tmp_path: Path):
    image_path = tmp_path / "logo.png"
    image_path.write_bytes(
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR"
        b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00"
        b"\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cc\xf8\xff\xff?"
        b"\x00\x05\xfe\x02\xfeA\x88\x97\x19\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    pptx_path = tmp_path / "with_logo.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(0.2), Inches(0.2), Inches(0.8), Inches(0.8))
    prs.save(str(pptx_path))

    context = analyze_reference_pptx(pptx_path, asset_dir=tmp_path / "assets")

    assert context.logo_path
    assert Path(context.logo_path).exists()
    assert context.image_paths == [context.logo_path]

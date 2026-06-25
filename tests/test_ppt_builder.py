from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.models import GeneratedSlideContent, PPTGenerationInput, ReferenceContext
from app.ppt_builder import build_pptx


@pytest.fixture
def sample_input() -> PPTGenerationInput:
    return PPTGenerationInput(
        topic="测试主题",
        purpose="内部培训",
        audience="开发人员",
        slide_count=3,
        style="专业正式",
        extra_notes="",
    )


@pytest.fixture
def sample_slides() -> list[GeneratedSlideContent]:
    return [
        GeneratedSlideContent(
            slide_index=1,
            title="第一页",
            bullets=["要点一", "要点二", "要点三"],
            speaker_hint="开场介绍",
        ),
        GeneratedSlideContent(
            slide_index=2,
            title="第二页",
            bullets=["核心观点", "数据支撑", "案例分析"],
            speaker_hint="讲解核心",
        ),
        GeneratedSlideContent(
            slide_index=3,
            title="第三页",
            bullets=["总结"],
            speaker_hint="收尾",
        ),
    ]


def test_build_pptx_creates_valid_pptx_file(
    tmp_path: Path, sample_input, sample_slides
):
    output_path = tmp_path / "test_output.pptx"
    result = build_pptx(sample_input, sample_slides, output_path)
    assert result.exists()
    assert result.suffix == ".pptx"


def test_build_pptx_has_correct_number_of_slides(
    tmp_path: Path, sample_input, sample_slides
):
    output_path = tmp_path / "test.pptx"
    build_pptx(sample_input, sample_slides, output_path)
    prs = Presentation(str(output_path))
    # 1 title slide + 3 content slides = 4
    assert len(prs.slides) == 4


def test_build_pptx_title_slide_contains_topic(
    tmp_path: Path, sample_input, sample_slides
):
    output_path = tmp_path / "test.pptx"
    build_pptx(sample_input, sample_slides, output_path)
    prs = Presentation(str(output_path))
    first_slide = prs.slides[0]
    # 检查第一页（封面）的文本
    all_text = ""
    for shape in first_slide.shapes:
        if hasattr(shape, "text"):
            all_text += shape.text
    assert "测试主题" in all_text


def test_build_pptx_content_slides_contain_titles(
    tmp_path: Path, sample_input, sample_slides
):
    output_path = tmp_path / "test.pptx"
    build_pptx(sample_input, sample_slides, output_path)
    prs = Presentation(str(output_path))
    # 从第 2 页开始是正文
    for i, content in enumerate(sample_slides):
        slide = prs.slides[i + 1]
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text
        assert content.title in slide_text


def test_build_pptx_content_slides_contain_bullets(
    tmp_path: Path, sample_input, sample_slides
):
    output_path = tmp_path / "test.pptx"
    build_pptx(sample_input, sample_slides, output_path)
    prs = Presentation(str(output_path))
    slide = prs.slides[1]  # 第一页正文
    slide_text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            slide_text += shape.text
    assert "要点一" in slide_text


def test_build_pptx_raises_on_empty_slides(tmp_path: Path, sample_input):
    with pytest.raises(ValueError, match="没有可生成"):
        build_pptx(sample_input, [], tmp_path / "empty.pptx")


def test_build_pptx_reuses_reference_logo_on_slides(
    tmp_path: Path, sample_input, sample_slides
):
    logo_path = tmp_path / "logo.png"
    logo_path.write_bytes(
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR"
        b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00"
        b"\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cc\xf8\xff\xff?"
        b"\x00\x05\xfe\x02\xfeA\x88\x97\x19\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    ref_context = ReferenceContext(
        slide_count=1,
        slides_summary=[],
        logo_path=str(logo_path),
        image_paths=[str(logo_path)],
    )
    output_path = tmp_path / "with_logo.pptx"

    build_pptx(sample_input, sample_slides, output_path, ref_context)

    prs = Presentation(str(output_path))
    assert any(shape.shape_type == 13 for shape in prs.slides[0].shapes)
    assert any(shape.shape_type == 13 for shape in prs.slides[1].shapes)
